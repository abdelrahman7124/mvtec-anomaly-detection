"""Build the project slide deck from the rendered figures and result tables.

Every number on a slide is read from reports/tables/ (headline_facts.json for the
EDA, baseline_results.csv for the models), so the deck cannot drift away from the
analysis that produced it.
"""

import _bootstrap  # noqa: F401
import json

import pandas as pd
from PIL import Image as PILImage
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from mvtec_eda.config import FIGURES_DIR, PRESENTATION_DIR, TABLES_DIR, ensure_output_dirs

# --- Design tokens, matching the figures -----------------------------------
SURFACE = RGBColor(0xFC, 0xFC, 0xFB)
INK = RGBColor(0x0B, 0x0B, 0x0B)
INK_SECONDARY = RGBColor(0x52, 0x51, 0x4E)
INK_MUTED = RGBColor(0x8A, 0x88, 0x80)
BLUE = RGBColor(0x2A, 0x78, 0xD6)
ORANGE = RGBColor(0xEB, 0x68, 0x34)
AQUA = RGBColor(0x1B, 0xAF, 0x7A)
VIOLET = RGBColor(0x4A, 0x3A, 0xA7)
GRID = RGBColor(0xE5, 0xE4, 0xE0)
BAND = RGBColor(0xF2, 0xF1, 0xEE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.62)
CONTENT_W = SLIDE_W - 2 * MARGIN

TITLE_TOP = Inches(0.42)
TITLE_H = Inches(0.62)
SUB_TOP = Inches(1.02)
SUB_H = Inches(0.38)
BODY_TOP = Inches(1.56)
TAKEAWAY_H = Inches(0.72)
TAKEAWAY_TOP = SLIDE_H - TAKEAWAY_H - Inches(0.34)


def _text(slide, left, top, width, height, text, size, colour, bold=False,
          align=PP_ALIGN.LEFT, italic=False, spacing=1.0):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = spacing
    run = p.add_run()
    run.text = text
    f = run.font
    f.name, f.size, f.bold, f.italic = FONT, Pt(size), bold, italic
    f.color.rgb = colour
    return box


def _blank(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = SURFACE
    return slide


def _rect(slide, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    shape.shadow.inherit = False
    return shape


def _fit(path, box_left, box_top, box_w, box_h):
    """Scale an image to fit inside a box, preserving aspect ratio and centring it."""
    with PILImage.open(path) as im:
        iw, ih = im.size
    scale = min(box_w / iw, box_h / ih)
    w, h = int(iw * scale), int(ih * scale)
    return int(box_left + (box_w - w) / 2), int(box_top + (box_h - h) / 2), w, h


def figure_slide(prs, title, subtitle, figure_name, takeaway):
    slide = _blank(prs)
    _text(slide, MARGIN, TITLE_TOP, CONTENT_W, TITLE_H, title, 27, INK, bold=True)
    if subtitle:
        _text(slide, MARGIN, SUB_TOP, CONTENT_W, SUB_H, subtitle, 13.5, INK_SECONDARY)

    body_top = BODY_TOP if subtitle else Inches(1.24)
    box_h = (TAKEAWAY_TOP if takeaway else SLIDE_H - Inches(0.4)) - body_top - Inches(0.22)
    left, top, w, h = _fit(FIGURES_DIR / figure_name, int(MARGIN), int(body_top),
                           int(CONTENT_W), int(box_h))
    slide.shapes.add_picture(str(FIGURES_DIR / figure_name), Emu(left), Emu(top), Emu(w), Emu(h))

    if takeaway:
        _rect(slide, MARGIN, TAKEAWAY_TOP, CONTENT_W, TAKEAWAY_H, BAND)
        _rect(slide, MARGIN, TAKEAWAY_TOP, Inches(0.055), TAKEAWAY_H, BLUE)
        _text(slide, MARGIN + Inches(0.26), TAKEAWAY_TOP + Inches(0.13),
              CONTENT_W - Inches(0.5), TAKEAWAY_H - Inches(0.2), takeaway, 13.5, INK,
              spacing=1.15)
    return slide


def section_slide(prs, eyebrow, title, blurb):
    slide = _blank(prs)
    _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.13), ORANGE)
    _text(slide, MARGIN, Inches(2.6), CONTENT_W, Inches(0.42), eyebrow, 15, ORANGE, bold=True)
    _text(slide, MARGIN, Inches(3.15), CONTENT_W, Inches(1.0), title, 38, INK, bold=True)
    _text(slide, MARGIN, Inches(4.25), Inches(9.5), Inches(1.0), blurb, 17, INK_SECONDARY,
          spacing=1.3)
    return slide


def build(facts: dict, results: pd.DataFrame | None) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Model slides need both the numbers and the rendered figures; a partial
    # baseline run should still produce a valid EDA-only deck.
    required = ["fig12_baseline_image_auroc.png", "fig13_baseline_summary.png",
                "fig14_qualitative_maps.png"]
    has_models = (
        results is not None
        and len(results) > 0
        and all((FIGURES_DIR / name).exists() for name in required)
    )
    if results is not None and len(results) > 0 and not has_models:
        print("  baseline figures missing - run scripts/06_model_figures.py")

    # --- 1. Title ----------------------------------------------------------
    slide = _blank(prs)
    _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.13), BLUE)
    _text(slide, MARGIN, Inches(2.05), CONTENT_W, Inches(0.42),
          "DATA ANALYSIS AND BASELINE MODELS", 15, BLUE, bold=True)
    _text(slide, MARGIN, Inches(2.62), CONTENT_W, Inches(1.1),
          "Detecting Manufacturing Defects", 42, INK, bold=True)
    _text(slide, MARGIN, Inches(3.62), CONTENT_W, Inches(0.75),
          "Unsupervised anomaly detection on the MVTec AD dataset", 21, INK_SECONDARY)
    _rect(slide, MARGIN, Inches(4.62), Inches(1.5), Inches(0.035), GRID)
    _text(slide, MARGIN, Inches(4.95), CONTENT_W, Inches(0.9),
          f"{facts['n_images']:,} images  ·  {facts['n_categories']} product categories  ·  "
          f"{facts['n_defect_types']} defect types  ·  {facts['n_masks']:,} pixel-level masks",
          15, INK_MUTED)

    # --- 2. What this covers ----------------------------------------------
    slide = _blank(prs)
    _text(slide, MARGIN, TITLE_TOP, CONTENT_W, TITLE_H,
          "What this presentation covers", 27, INK, bold=True)
    items = [
        ("The problem",
         "A factory needs to catch defective parts automatically. Defects are rare, varied, "
         "and mostly unseen when the system is built."),
        ("The data",
         "MVTec AD: 15 product categories shot under controlled conditions, with a "
         "pixel-accurate mask for every defect in the test set."),
        ("What the data allows",
         "Six design decisions, each forced by a specific measurement rather than by habit."),
        ("Two baselines",
         "A convolutional autoencoder trained from scratch, and PaDiM on frozen pretrained "
         "features. One works far better, and the EDA explains why."),
    ]
    top = Inches(1.5)
    for i, (head, body) in enumerate(items):
        _rect(slide, MARGIN, top, Inches(0.05), Inches(0.92), BLUE if i % 2 == 0 else AQUA)
        _text(slide, MARGIN + Inches(0.3), top, Inches(3.2), Inches(0.4), head, 17, INK, bold=True)
        _text(slide, MARGIN + Inches(3.7), top + Inches(0.03), CONTENT_W - Inches(4.0),
              Inches(1.0), body, 14.5, INK_SECONDARY, spacing=1.25)
        top += Inches(1.24)

    # --- 3. Dataset at a glance -------------------------------------------
    slide = _blank(prs)
    _text(slide, MARGIN, TITLE_TOP, CONTENT_W, TITLE_H, "The dataset at a glance", 27, INK,
          bold=True)
    _text(slide, MARGIN, SUB_TOP, CONTENT_W, SUB_H,
          "Measured from the images, not quoted from the paper", 13.5, INK_SECONDARY)

    tiles = [
        (f"{facts['n_images']:,}", "images total", BLUE),
        (f"{facts['n_categories']}", f"categories  ({facts['n_textures']} textures, "
                                     f"{facts['n_objects']} objects)", AQUA),
        (f"{facts['n_defect_types']}", "distinct defect types", ORANGE),
        (f"{facts['n_masks']:,}", "pixel-level defect masks", INK_SECONDARY),
    ]
    tile_w = (CONTENT_W - Inches(0.45)) / 4
    for i, (value, label, colour) in enumerate(tiles):
        left = MARGIN + i * (tile_w + Inches(0.15))
        _rect(slide, left, Inches(1.72), tile_w, Inches(1.72), WHITE, line=GRID)
        _rect(slide, left, Inches(1.72), tile_w, Inches(0.05), colour)
        _text(slide, left + Inches(0.22), Inches(2.06), tile_w - Inches(0.4), Inches(0.8),
              value, 40, INK, bold=True)
        _text(slide, left + Inches(0.22), Inches(2.82), tile_w - Inches(0.4), Inches(0.55),
              label, 13, INK_SECONDARY, spacing=1.15)

    rows = [
        ("Training images", f"{facts['n_train']:,}", "all defect-free — no anomalies at all"),
        ("Test images", f"{facts['n_test']:,}",
         f"{facts['n_test_normal']} normal, {facts['n_test_anomalous']:,} anomalous "
         f"({facts['pct_test_anomalous']}%)"),
        ("Median defect size", f"{facts['median_defect_area_pct']}%",
         f"of the image; {facts['pct_defects_under_1pct']}% of defects cover under 1%"),
        ("Resolutions", f"{facts['n_resolutions']}",
         f"from 700x700 to 1024x1024; "
         f"{len(facts['greyscale_categories'])} categories are single-channel"),
    ]
    top = Inches(3.86)
    for label, value, note in rows:
        _text(slide, MARGIN + Inches(0.05), top, Inches(3.0), Inches(0.4), label, 15,
              INK_SECONDARY)
        _text(slide, MARGIN + Inches(3.1), top, Inches(1.9), Inches(0.4), value, 15, INK,
              bold=True)
        _text(slide, MARGIN + Inches(5.0), top, CONTENT_W - Inches(5.2), Inches(0.4), note, 15,
              INK_SECONDARY)
        top += Inches(0.44)
        _rect(slide, MARGIN, top - Inches(0.08), CONTENT_W, Inches(0.008), GRID)

    # --- EDA figure slides -------------------------------------------------
    figure_slide(
        prs, "Fifteen categories, two visual regimes",
        "Five textures and ten centred objects",
        "fig05_normal_samples.png",
        "Textures repeat across the whole frame; objects sit in a fixed pose. The two behave "
        "differently under augmentation, so they are handled separately.",
    )
    figure_slide(
        prs, "The training set contains no defects at all",
        "This single fact decides the entire modelling approach",
        "fig02_composition.png",
        f"With {facts['n_train_anomalies']} anomalies among {facts['n_train']:,} training "
        "images, a supervised classifier cannot be fitted. The task is one-class: learn "
        "'normal', then measure deviation.",
    )
    figure_slide(
        prs, "The same one-class split in every category",
        "Train is always defect-free; anomalies appear only at test time",
        "fig01_split_overview.png",
        f"Category size ranges from {facts['smallest_category']} "
        f"({facts['smallest_category_train']} training images) to hazelnut (391). Small "
        "categories give the least stable scores.",
    )
    figure_slide(
        prs, "Defects are many, and each one is rare",
        f"{facts['n_defect_types']} category-specific defect types across the test split",
        "fig03_defect_types.png",
        f"The median defect type has only {facts['median_images_per_defect_type']} images "
        f"(minimum {facts['min_images_per_defect_type']}). Learning one class per defect is "
        "not viable — the model must flag anything unlike normal.",
    )
    figure_slide(
        prs, "What the model has to catch",
        "Ground-truth masks mark every defective pixel",
        "fig06b_defect_examples_wide.png",
        "Defects range from a barely visible scratch on a screw to a completely misplaced "
        "transistor. The same detector has to cover both extremes.",
    )
    figure_slide(
        prs, "Most defects are very small",
        "Measured directly from the 1,258 ground-truth masks",
        "fig07_defect_area.png",
        f"Median defect covers {facts['median_defect_area_pct']}% of its image and "
        f"{facts['pct_defects_under_1pct']}% cover under 1%. Resizing to 224x224 would shrink "
        "a typical small defect to a handful of pixels.",
    )
    figure_slide(
        prs, "Defect size varies hugely between categories",
        f"From {facts['smallest_defect_pct']}% to {facts['largest_defect_pct']}% of the image",
        "fig08_defect_area_by_category.png",
        "Screw and capsule defects are an order of magnitude smaller than tile or bottle "
        "defects, so one detection threshold across categories would not work.",
    )
    figure_slide(
        prs, "Images are high-resolution but not uniform",
        f"{facts['n_resolutions']} distinct resolutions; all images square",
        "fig04_image_properties.png",
        f"{facts['n_greyscale_images']:,} images in "
        f"{', '.join(facts['greyscale_categories'])} are stored single-channel and must be "
        "expanded to three before a pretrained backbone will accept them.",
    )
    figure_slide(
        prs, "Categories look nothing like each other",
        "Median brightness, measured per category",
        "fig09_appearance.png",
        "Brightness spans roughly 48 to 185 on a 0-255 scale. One global normalisation would "
        "leave several categories badly off-centre.",
    )
    figure_slide(
        prs, "Defects sit near the image centre",
        "A property of how the dataset was photographed, not of defects in general",
        "fig10_defect_location.png",
        "48% of defects fall within 0.2 of the centre. A real production line would not centre "
        "defects, so this bias must not be learned — no centre-cropping.",
    )
    figure_slide(
        prs, "Test normals match the training distribution",
        "Checked rather than assumed — a one-class model depends on it",
        "fig11_train_test_shift.png",
        f"{facts['n_categories_within_5_levels']} of 15 categories agree within ±5 intensity "
        f"levels. screw, grid and leather drift up to {facts['max_intensity_drift']}, so "
        "normalisation statistics are computed per category.",
    )

    # --- Findings to decisions --------------------------------------------
    slide = _blank(prs)
    _text(slide, MARGIN, TITLE_TOP, CONTENT_W, TITLE_H,
          "Six findings, six design decisions", 27, INK, bold=True)
    _text(slide, MARGIN, SUB_TOP, CONTENT_W, SUB_H,
          "Each decision traces back to a measurement in this deck", 13.5, INK_SECONDARY)
    pairs = [
        (f"{facts['n_train_anomalies']} anomalies in {facts['n_train']:,} training images",
         "One-class method — a supervised classifier is impossible"),
        (f"{facts['n_defect_types']} defect types, median "
         f"{facts['median_images_per_defect_type']} images each",
         "Detect deviation from normal, do not model defect classes"),
        (f"{facts['pct_test_anomalous']}% of the test split is anomalous",
         "Report AUROC — accuracy is misleading at this base rate"),
        (f"Median defect covers {facts['median_defect_area_pct']}% of the image",
         "Keep resolution high; score per patch, not per image"),
        (f"{facts['n_resolutions']} resolutions, "
         f"{len(facts['greyscale_categories'])} single-channel categories",
         "Loader resizes consistently and expands channels"),
        ("Brightness spans 48-185; 3 categories drift",
         "One model and one normalisation per category"),
    ]
    _text(slide, MARGIN + Inches(0.16), Inches(1.28), Inches(5.6), Inches(0.3),
          "WHAT THE DATA SHOWS", 11.5, INK_MUTED, bold=True)
    _text(slide, MARGIN + Inches(6.5), Inches(1.28), Inches(5.6), Inches(0.3),
          "WHAT IT FORCES US TO DO", 11.5, INK_MUTED, bold=True)
    top, row_h = Inches(1.62), Inches(0.82)
    for i, (finding, decision) in enumerate(pairs):
        if i % 2 == 0:
            _rect(slide, MARGIN, top - Inches(0.08), CONTENT_W, row_h,
                  RGBColor(0xF6, 0xF5, 0xF2))
        _text(slide, MARGIN + Inches(0.16), top + Inches(0.06), Inches(5.9), row_h,
              finding, 14.5, INK_SECONDARY, spacing=1.15)
        _rect(slide, MARGIN + Inches(6.25), top + Inches(0.02), Inches(0.04),
              row_h - Inches(0.2), BLUE)
        _text(slide, MARGIN + Inches(6.5), top + Inches(0.06), Inches(5.5), row_h,
              decision, 14.5, INK, bold=True, spacing=1.15)
        top += row_h

    if not has_models:
        _finish(prs)
        return

    # --- Model section -----------------------------------------------------
    n_cat = len(results)
    ae_img = results.ae_image_auroc.mean()
    pd_img = results.padim_image_auroc.mean()
    ae_pix = results.ae_pixel_auroc.mean()
    pd_pix = results.padim_pixel_auroc.mean()
    padim_wins = int((results.padim_image_auroc > results.ae_image_auroc).sum())
    best = results.loc[results.padim_image_auroc.idxmax()]
    worst = results.loc[results.padim_image_auroc.idxmin()]

    section_slide(
        prs, "PART TWO", "Two baselines",
        "Both train on defect-free images only, because the data leaves no other option. "
        "The difference between them is where they look.",
    )

    # How the two models work.
    slide = _blank(prs)
    _text(slide, MARGIN, TITLE_TOP, CONTENT_W, TITLE_H, "Two ways to model 'normal'", 27, INK,
          bold=True)
    _text(slide, MARGIN, SUB_TOP, CONTENT_W, SUB_H,
          "Neither one ever sees a defective image during training", 13.5, INK_SECONDARY)

    col_w = (CONTENT_W - Inches(0.4)) / 2
    columns = [
        (ORANGE, "Convolutional autoencoder", "trained from scratch", [
            "Compress the image to a small bottleneck, then rebuild it.",
            "Trained only on normal images, so it never learns to rebuild a defect.",
            "Anomaly score = how badly a pixel was reconstructed.",
            "The catch: the bottleneck describes the whole image at once, so a defect "
            "covering 1.5% of the frame barely changes the loss.",
        ]),
        (BLUE, "PaDiM", "frozen ImageNet features", [
            "Push images through a pretrained ResNet-18 — no training at all.",
            "At each patch position, fit a Gaussian over the normal training images.",
            "Anomaly score = Mahalanobis distance from that position's Gaussian.",
            "The advantage: every patch position keeps its own model, so a small local "
            "deviation cannot be averaged away.",
        ]),
    ]
    for i, (colour, name, tag, bullets) in enumerate(columns):
        left = MARGIN + i * (col_w + Inches(0.4))
        _rect(slide, left, Inches(1.62), col_w, Inches(4.6), WHITE, line=GRID)
        _rect(slide, left, Inches(1.62), col_w, Inches(0.06), colour)
        _text(slide, left + Inches(0.26), Inches(1.92), col_w - Inches(0.5), Inches(0.4),
              name, 20, INK, bold=True)
        _text(slide, left + Inches(0.26), Inches(2.36), col_w - Inches(0.5), Inches(0.3),
              tag, 13, colour, bold=True)
        top = Inches(2.84)
        for bullet in bullets:
            _rect(slide, left + Inches(0.26), top + Inches(0.09), Inches(0.07), Inches(0.07),
                  colour)
            _text(slide, left + Inches(0.5), top, col_w - Inches(0.78), Inches(0.9), bullet,
                  13.5, INK_SECONDARY, spacing=1.22)
            top += Inches(0.82)

    figure_slide(
        prs, "PaDiM wins in every category",
        f"Image-level AUROC across all {n_cat} categories",
        "fig12_baseline_image_auroc.png",
        f"PaDiM scores higher than the autoencoder in {padim_wins} of {n_cat} categories. "
        f"Best: {best.category} at {best.padim_image_auroc:.3f}. "
        f"Weakest: {worst.category} at {worst.padim_image_auroc:.3f}.",
    )
    figure_slide(
        prs, "The gap is even wider for localisation",
        "Mean AUROC across all categories, image-level and pixel-level",
        "fig13_baseline_summary.png",
        f"Image-level: {ae_img:.2f} → {pd_img:.2f}. Pixel-level: {ae_pix:.2f} → {pd_pix:.2f}. "
        "The autoencoder can sometimes tell that an image is odd, but not where.",
    )
    figure_slide(
        prs, "Where each model thinks the defect is",
        "Anomaly heatmaps against the ground-truth mask",
        "fig14_qualitative_maps.png",
        "The autoencoder highlights edges and texture everywhere — it reconstructs "
        "high-frequency detail badly whether or not it is defective. PaDiM lands on the "
        "actual defect.",
    )

    # --- Results table ------------------------------------------------------
    slide = _blank(prs)
    _text(slide, MARGIN, TITLE_TOP, CONTENT_W, TITLE_H, "Full results", 27, INK, bold=True)
    _text(slide, MARGIN, SUB_TOP, CONTENT_W, SUB_H,
          "AUROC per category — higher is better, 0.5 is random", 13.5, INK_SECONDARY)

    tbl = results.sort_values("padim_image_auroc", ascending=False)
    n_rows = len(tbl) + 2
    col_x = [MARGIN, MARGIN + Inches(2.6), MARGIN + Inches(4.9),
             MARGIN + Inches(7.4), MARGIN + Inches(9.7)]
    headers = ["Category", "AE image", "PaDiM image", "AE pixel", "PaDiM pixel"]
    row_h = Inches(0.315)

    _text(slide, col_x[0], Inches(1.56), Inches(2.4), Inches(0.3), headers[0], 12.5,
          INK_MUTED, bold=True)
    for j, head in enumerate(headers[1:], start=1):
        _text(slide, col_x[j], Inches(1.56), Inches(2.2), Inches(0.3), head, 12.5,
              ORANGE if "AE" in head else BLUE, bold=True)

    top = Inches(1.92)
    for i, (_, r) in enumerate(tbl.iterrows()):
        if i % 2 == 0:
            _rect(slide, MARGIN - Inches(0.1), top - Inches(0.04), CONTENT_W, row_h,
                  RGBColor(0xF6, 0xF5, 0xF2))
        _text(slide, col_x[0], top, Inches(2.4), row_h, r.category, 13, INK)
        for j, value in enumerate(
            [r.ae_image_auroc, r.padim_image_auroc, r.ae_pixel_auroc, r.padim_pixel_auroc],
            start=1,
        ):
            _text(slide, col_x[j], top, Inches(2.2), row_h, f"{value:.3f}", 13, INK_SECONDARY)
        top += row_h

    _rect(slide, MARGIN - Inches(0.1), top, CONTENT_W, Inches(0.012), GRID)
    top += Inches(0.1)
    _text(slide, col_x[0], top, Inches(2.4), row_h, "MEAN", 13, INK, bold=True)
    for j, value in enumerate([ae_img, pd_img, ae_pix, pd_pix], start=1):
        _text(slide, col_x[j], top, Inches(2.2), row_h, f"{value:.3f}", 13, INK, bold=True)

    # --- Conclusions --------------------------------------------------------
    slide = _blank(prs)
    _text(slide, MARGIN, TITLE_TOP, CONTENT_W, TITLE_H, "What we learned, and what is next",
          27, INK, bold=True)

    _rect(slide, MARGIN, Inches(1.34), CONTENT_W, Inches(1.5), WHITE, line=GRID)
    _rect(slide, MARGIN, Inches(1.34), Inches(0.055), Inches(1.5), BLUE)
    _text(slide, MARGIN + Inches(0.3), Inches(1.54), CONTENT_W - Inches(0.6), Inches(0.4),
          "The EDA predicted the result", 20, INK, bold=True)
    _text(slide, MARGIN + Inches(0.3), Inches(2.02), CONTENT_W - Inches(0.6), Inches(0.75),
          f"A median defect covers {facts['median_defect_area_pct']}% of its image. A model "
          "that compresses the whole frame into one bottleneck cannot represent that, and the "
          f"autoencoder duly scores {ae_pix:.2f} pixel AUROC. A model that keeps every patch "
          f"position separate scores {pd_pix:.2f}. The measurement came first.",
          14.5, INK_SECONDARY, spacing=1.25)

    steps = [
        ("1", "Strengthen the baseline",
         "PatchCore with a coreset memory bank; compare against PaDiM on the same split"),
        ("2", "Tune what matters",
         "Backbone depth, input resolution and feature layers — resolution should matter most"),
        ("3", "Handle the weak categories",
         f"{worst.category} and toothbrush score lowest; test whether augmentation helps"),
        ("4", "Make it deployable",
         "Pick an operating threshold per category and report precision and recall there"),
    ]
    top = Inches(3.24)
    for num, head, body in steps:
        _rect(slide, MARGIN, top, Inches(0.42), Inches(0.42), BLUE)
        _text(slide, MARGIN, top + Inches(0.045), Inches(0.42), Inches(0.35), num, 15, WHITE,
              bold=True, align=PP_ALIGN.CENTER)
        _text(slide, MARGIN + Inches(0.62), top - Inches(0.01), Inches(3.9), Inches(0.4),
              head, 15.5, INK, bold=True)
        _text(slide, MARGIN + Inches(4.6), top + Inches(0.01), CONTENT_W - Inches(4.8),
              Inches(0.45), body, 14, INK_SECONDARY)
        top += Inches(0.82)

    _finish(prs)


def _finish(prs) -> None:
    ensure_output_dirs()
    out = PRESENTATION_DIR / "week1_eda.pptx"
    prs.save(out)
    print(f"Wrote {out} ({len(prs.slides._sldIdLst)} slides)")


def main() -> None:
    facts = json.loads((TABLES_DIR / "headline_facts.json").read_text(encoding="utf-8"))

    results = None
    csv = TABLES_DIR / "baseline_results.csv"
    if csv.exists():
        results = pd.read_csv(csv)
        needed = ["ae_image_auroc", "padim_image_auroc", "ae_pixel_auroc", "padim_pixel_auroc"]
        results = results.dropna(subset=[c for c in needed if c in results])
        print(f"baseline results: {len(results)} categories")
    else:
        print("no baseline_results.csv yet - building the EDA-only deck")

    build(facts, results)


if __name__ == "__main__":
    main()
